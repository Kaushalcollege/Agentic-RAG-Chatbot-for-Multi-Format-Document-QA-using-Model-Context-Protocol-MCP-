# backend/utils/parser.py
import fitz  # PyMuPDF
import docx
import csv
import pptx
import pandas as pd
from pptx import Presentation
import io

def parse_pdf(file_bytes):
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def parse_docx(file_bytes: bytes) -> str:
    doc_stream = io.BytesIO(file_bytes)  # Wrap bytes as file-like
    doc = docx.Document(doc_stream)      # Works now
    return "\n".join([para.text for para in doc.paragraphs])


def parse_pptx(file_bytes: bytes) -> str:
    # Convert bytes to a file-like object
    pptx_stream = io.BytesIO(file_bytes)
    prs = Presentation(pptx_stream)

    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    return "\n".join(text_runs)

def parse_csv(file_bytes: bytes) -> str:
    csv_stream = io.BytesIO(file_bytes)  # wrap in file-like object
    df = pd.read_csv(csv_stream)
    return df.to_string(index=False)

def parse_txt(file_bytes):
    return file_bytes.read().decode("utf-8")

def parse_file(filename, file_bytes):
    if filename.endswith(".pdf"):
        return parse_pdf(file_bytes)
    elif filename.endswith(".docx"):
        return parse_docx(file_bytes)
    elif filename.endswith(".pptx"):
        return parse_pptx(file_bytes)
    elif filename.endswith(".csv"):
        return parse_csv(file_bytes)
    elif filename.endswith(".txt") or filename.endswith(".md"):
        return parse_txt(file_bytes)
    else:
        return ""
